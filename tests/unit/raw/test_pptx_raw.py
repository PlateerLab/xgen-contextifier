# tests/unit/raw/test_pptx_raw.py
"""PptxRawDocument — C4 of the raw layer.

Decks are built with python-pptx (the reference implementation) and
reopened with it after every raw edit, so each assertion is a
round-trip through an independent OOXML reader.
"""

from __future__ import annotations

import base64
import io
import zipfile

import pytest

pptx = pytest.importorskip("pptx")

from pptx import Presentation  # noqa: E402
from pptx.chart.data import CategoryChartData  # noqa: E402
from pptx.enum.chart import XL_CHART_TYPE  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from pptx.util import Inches  # noqa: E402

from xgen_contextifier.raw import open_raw  # noqa: E402

# 1x1 px valid PNG — fallback when Pillow is unavailable.
_TINY_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR4nGP4"
    "z8DwHwAFBQIAX8jx0gAAAABJRU5ErkJggg=="
)


def _png_bytes() -> bytes:
    try:
        from PIL import Image
    except ImportError:
        return _TINY_PNG
    buf = io.BytesIO()
    Image.new("RGB", (8, 8), (200, 30, 30)).save(buf, format="PNG")
    return buf.getvalue()


def build_deck() -> bytes:
    """Slide 1: title (bold run) + body + notes. Slide 2: table, chart,
    picture."""
    prs = Presentation()

    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Hello Title"
    s1.shapes.title.text_frame.paragraphs[0].runs[0].font.bold = True
    s1.placeholders[1].text = "Body text"
    s1.notes_slide.notes_text_frame.text = "Speaker notes here"

    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = s2.shapes.add_table(2, 3, Inches(0.5), Inches(0.5), Inches(6), Inches(1.5))
    for r in range(2):
        for c in range(3):
            tbl.table.cell(r, c).text = f"R{r}C{c}"
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Sales", (10.0, 20.0, 30.0))
    s2.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.5),
        Inches(2.5),
        Inches(4),
        Inches(3),
        chart_data,
    )
    s2.shapes.add_picture(
        io.BytesIO(_png_bytes()), Inches(5), Inches(2.5), Inches(1), Inches(1)
    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_three_slide_deck() -> bytes:
    """Three text slides; slide 2 additionally owns the deck's ONLY
    chart (with embedded workbook) and a notes slide."""
    prs = Presentation()
    for i, txt in enumerate(["One", "Two", "Three"]):
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
        box.text_frame.text = txt
        if i == 1:
            chart_data = CategoryChartData()
            chart_data.categories = ["A", "B"]
            chart_data.add_series("S", (1.0, 2.0))
            s.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED,
                Inches(1),
                Inches(2),
                Inches(3),
                Inches(3),
                chart_data,
            )
            s.notes_slide.notes_text_frame.text = "gone with the slide"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# A complete handwritten slide: one text box saying REPLACED.
MINIMAL_SLIDE_XML = (
    b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    b'<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'
    b' xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
    b' xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
    b"<p:cSld><p:spTree>"
    b'<p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr>'
    b"<p:grpSpPr/>"
    b"<p:sp><p:nvSpPr>"
    b'<p:cNvPr id="2" name="TextBox 1"/><p:cNvSpPr txBox="1"/><p:nvPr/>'
    b"</p:nvSpPr>"
    b'<p:spPr><a:xfrm><a:off x="914400" y="914400"/>'
    b'<a:ext cx="3657600" cy="457200"/></a:xfrm>'
    b'<a:prstGeom prst="rect"><a:avLst/></a:prstGeom></p:spPr>'
    b"<p:txBody><a:bodyPr/><a:lstStyle/>"
    b"<a:p><a:r><a:t>REPLACED</a:t></a:r></a:p></p:txBody>"
    b"</p:sp>"
    b"</p:spTree></p:cSld>"
    b"<p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>"
    b"</p:sld>"
)


class TestInventory:
    def test_shapes_kinds_and_texts(self):
        data = build_deck()
        raw = open_raw(data)
        assert raw.format == "pptx"
        assert len(raw.slides) == 2
        assert raw.slides[0].part_name == "ppt/slides/slide1.xml"

        s1 = raw.slides[0].shapes
        assert [s.kind for s in s1] == ["text", "text"]
        assert [s.text for s in s1] == ["Hello Title", "Body text"]

        s2 = raw.slides[1].shapes
        assert [s.kind for s in s2] == ["table", "chart", "picture"]
        assert all(s.text is None for s in s2)
        assert all(s.name for s in s2)  # python-pptx names every shape

    def test_slide_and_shape_order_matches_python_pptx(self):
        data = build_deck()
        raw = open_raw(data)
        prs = Presentation(io.BytesIO(data))
        assert len(raw.slides) == len(prs.slides)
        for raw_slide, pp_slide in zip(raw.slides, prs.slides):
            assert [s.id for s in raw_slide.shapes] == [
                sh.shape_id for sh in pp_slide.shapes
            ]


class TestText:
    def test_set_text_roundtrip_preserves_first_run_formatting(self):
        raw = open_raw(build_deck())
        slide = raw.slides[0]
        title_id = next(s.id for s in slide.shapes if s.text == "Hello Title")
        slide.set_text(title_id, "New Title")
        assert slide.get_text(title_id) == "New Title"

        prs = Presentation(io.BytesIO(raw.to_bytes()))
        title = prs.slides[0].shapes.title
        assert title.text == "New Title"
        run = title.text_frame.paragraphs[0].runs[0]
        assert run.font.bold is True  # a:rPr of the original run survived

    def test_get_and_set_text_reject_shapes_without_txbody(self):
        raw = open_raw(build_deck())
        slide = raw.slides[1]
        pic_id = next(s.id for s in slide.shapes if s.kind == "picture")
        with pytest.raises(ValueError, match="text body"):
            slide.get_text(pic_id)
        with pytest.raises(ValueError, match="text body"):
            slide.set_text(pic_id, "nope")

    def test_set_text_paragraph_out_of_range(self):
        raw = open_raw(build_deck())
        slide = raw.slides[0]
        title_id = next(s.id for s in slide.shapes if s.text == "Hello Title")
        with pytest.raises(IndexError):
            slide.set_text(title_id, "x", para=5)


class TestTables:
    def test_dims_and_cell_read(self):
        raw = open_raw(build_deck())
        tables = raw.slides[1].tables
        assert len(tables) == 1
        t = tables[0]
        assert (t.n_rows, t.n_cols) == (2, 3)
        assert t.cell(0, 0).text == "R0C0"
        assert t.cell(1, 2).text == "R1C2"
        with pytest.raises(IndexError):
            t.cell(2, 0)

    def test_cell_set_text_and_insert_row_roundtrip(self):
        raw = open_raw(build_deck())
        t = raw.slides[1].tables[0]
        t.cell(0, 0).set_text("HDR")
        t.insert_row(1)  # empty row cloned from row 0
        assert (t.n_rows, t.n_cols) == (3, 3)

        prs = Presentation(io.BytesIO(raw.to_bytes()))
        table = next(sh for sh in prs.slides[1].shapes if sh.has_table).table
        assert len(table.rows) == 3
        assert len(table.columns) == 3
        assert table.cell(0, 0).text == "HDR"
        assert [table.cell(1, c).text for c in range(3)] == ["", "", ""]
        assert [table.cell(2, c).text for c in range(3)] == ["R1C0", "R1C1", "R1C2"]

    def test_delete_row_roundtrip(self):
        raw = open_raw(build_deck())
        raw.slides[1].tables[0].delete_row(0)
        prs = Presentation(io.BytesIO(raw.to_bytes()))
        table = next(sh for sh in prs.slides[1].shapes if sh.has_table).table
        assert len(table.rows) == 1
        assert table.cell(0, 0).text == "R1C0"


class TestChartsAndNotes:
    def test_chart_part_names(self):
        raw = open_raw(build_deck())
        assert raw.slides[1].chart_part_names == ["ppt/charts/chart1.xml"]
        assert raw.slides[0].chart_part_names == []

    def test_notes_text(self):
        raw = open_raw(build_deck())
        assert raw.slides[0].notes_text == "Speaker notes here"
        assert raw.slides[1].notes_text is None


class TestReplaceContent:
    def test_preserve_native_keeps_chart_table_picture(self):
        raw = open_raw(build_deck())
        preserved = raw.slides[1].replace_content(
            MINIMAL_SLIDE_XML, preserve_native=True
        )
        assert "table" in preserved
        assert "chart:chart1.xml" in preserved
        assert any(p.startswith("picture:") for p in preserved)

        prs = Presentation(io.BytesIO(raw.to_bytes()))
        slide = prs.slides[1]
        texts = [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame]
        assert "REPLACED" in texts

        chart_shapes = [sh for sh in slide.shapes if sh.has_chart]
        assert len(chart_shapes) == 1
        chart = chart_shapes[0].chart  # the native chart still loads...
        assert len(chart.plots) == 1  # ...and its plot area is readable
        assert list(chart.plots[0].categories) == ["Q1", "Q2", "Q3"]

        table_shapes = [sh for sh in slide.shapes if sh.has_table]
        assert len(table_shapes) == 1
        assert table_shapes[0].table.cell(0, 0).text == "R0C0"

        pictures = [
            sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        assert len(pictures) == 1

        ids = [sh.shape_id for sh in slide.shapes]
        assert len(ids) == len(set(ids))  # renumbering avoided collisions

    def test_no_preserve_drops_frames_but_keeps_parts(self):
        raw = open_raw(build_deck())
        preserved = raw.slides[1].replace_content(
            MINIMAL_SLIDE_XML, preserve_native=False
        )
        assert preserved == []

        out = raw.to_bytes()
        prs = Presentation(io.BytesIO(out))
        slide = prs.slides[1]
        assert not any(sh.has_chart for sh in slide.shapes)
        assert not any(sh.has_table for sh in slide.shapes)
        assert [sh.text_frame.text for sh in slide.shapes if sh.has_text_frame] == [
            "REPLACED"
        ]
        # No orphan cleanup here — that is remove_slide's job.
        names = zipfile.ZipFile(io.BytesIO(out)).namelist()
        assert "ppt/charts/chart1.xml" in names
        assert any(n.startswith("ppt/embeddings/") for n in names)

    def test_rejects_non_slide_xml(self):
        raw = open_raw(build_deck())
        with pytest.raises(ValueError, match="p:sld"):
            raw.slides[0].replace_content(b"<not-a-slide/>")


class TestRemoveSlide:
    def test_orphan_cleanup_and_byte_identity(self):
        data = build_three_slide_deck()
        zin = zipfile.ZipFile(io.BytesIO(data))
        assert "ppt/charts/chart1.xml" in zin.namelist()
        assert any(n.startswith("ppt/embeddings/") for n in zin.namelist())

        raw = open_raw(data)
        raw.remove_slide(1)
        out = raw.to_bytes()
        z = zipfile.ZipFile(io.BytesIO(out))
        names = set(z.namelist())

        # The slide, its notes, the chart and its embedded workbook are gone.
        for gone in (
            "ppt/slides/slide2.xml",
            "ppt/slides/_rels/slide2.xml.rels",
            "ppt/notesSlides/notesSlide1.xml",
            "ppt/notesSlides/_rels/notesSlide1.xml.rels",
            "ppt/charts/chart1.xml",
            "ppt/charts/_rels/chart1.xml.rels",
        ):
            assert gone not in names, f"orphan left behind: {gone}"
        assert not any(n.startswith("ppt/embeddings/") for n in names)

        # Shared infrastructure survives the reference count.
        assert "ppt/slideLayouts/slideLayout7.xml" in names
        assert "ppt/notesMasters/notesMaster1.xml" in names

        # Content types no longer mention the removed parts.
        content_types = z.read("[Content_Types].xml").decode()
        for part in (
            "/ppt/slides/slide2.xml",
            "/ppt/charts/chart1.xml",
            "/ppt/notesSlides/notesSlide1.xml",
        ):
            assert part not in content_types

        # Untouched slides are byte-identical.
        assert z.read("ppt/slides/slide1.xml") == zin.read("ppt/slides/slide1.xml")
        assert z.read("ppt/slides/slide3.xml") == zin.read("ppt/slides/slide3.xml")

        # And the deck still opens cleanly with the right slides.
        prs = Presentation(io.BytesIO(out))
        assert len(prs.slides) == 2
        remaining = [
            next(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
            for s in prs.slides
        ]
        assert remaining == ["One", "Three"]

    def test_index_out_of_range(self):
        raw = open_raw(build_three_slide_deck())
        with pytest.raises(IndexError):
            raw.remove_slide(3)


class TestBytePreservation:
    def test_set_text_touches_only_that_slide_part(self):
        data = build_deck()
        raw = open_raw(data)
        slide = raw.slides[0]
        title_id = next(s.id for s in slide.shapes if s.text == "Hello Title")
        slide.set_text(title_id, "Changed")
        out = raw.to_bytes()

        a = zipfile.ZipFile(io.BytesIO(data))
        b = zipfile.ZipFile(io.BytesIO(out))
        assert sorted(a.namelist()) == sorted(b.namelist())
        for name in a.namelist():
            if name == "ppt/slides/slide1.xml":
                assert a.read(name) != b.read(name)
            else:
                assert a.read(name) == b.read(name), f"collateral change in {name}"


def _slide_titles(data: bytes) -> list[str]:
    prs = Presentation(io.BytesIO(data))
    out = []
    for s in prs.slides:
        txt = ""
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                txt = sh.text_frame.text.strip().split("\n")[0]
                break
        out.append(txt)
    return out


class TestMoveSlide:
    def test_reorder_last_to_front(self):
        data = build_three_slide_deck()
        raw = open_raw(data)
        raw.move_slide(2, 0)
        out = raw.to_bytes()
        assert _slide_titles(out) == ["Three", "One", "Two"]

    def test_move_is_pure_presentation_reorder(self):
        """Only ppt/presentation.xml changes — every slide part byte-identical."""
        data = build_three_slide_deck()
        raw = open_raw(data)
        raw.move_slide(0, 2)
        out = raw.to_bytes()
        a = zipfile.ZipFile(io.BytesIO(data))
        b = zipfile.ZipFile(io.BytesIO(out))
        changed = [n for n in a.namelist() if a.read(n) != b.read(n)]
        assert changed == ["ppt/presentation.xml"]

    def test_out_of_range(self):
        raw = open_raw(build_three_slide_deck())
        with pytest.raises(IndexError):
            raw.move_slide(5, 0)
        with pytest.raises(IndexError):
            raw.move_slide(0, 9)

    def test_noop_same_position(self):
        data = build_three_slide_deck()
        raw = open_raw(data)
        raw.move_slide(1, 1)
        assert _slide_titles(raw.to_bytes()) == ["One", "Two", "Three"]


class TestDuplicateSlide:
    def test_duplicate_default_position_after_source(self):
        data = build_three_slide_deck()
        raw = open_raw(data)
        new_index = raw.duplicate_slide(0)
        out = raw.to_bytes()
        assert new_index == 1
        assert _slide_titles(out) == ["One", "One", "Two", "Three"]
        assert len(Presentation(io.BytesIO(out)).slides._sldIdLst) == 4

    def test_duplicate_at_explicit_position(self):
        raw = open_raw(build_three_slide_deck())
        raw.duplicate_slide(2, at=0)
        assert _slide_titles(raw.to_bytes())[0] == "Three"

    def test_duplicate_adds_parts_leaves_originals_byte_identical(self):
        """A duplicate ADDS parts; no pre-existing slide/chart part changes
        except presentation.xml + its rels + [Content_Types].xml."""
        data = build_three_slide_deck()  # slide 2 owns the only chart + notes
        raw = open_raw(data)
        raw.duplicate_slide(1)  # duplicate the chart+notes slide
        out = raw.to_bytes()
        a = zipfile.ZipFile(io.BytesIO(data))
        b = zipfile.ZipFile(io.BytesIO(out))
        added = sorted(n for n in b.namelist() if n not in a.namelist())
        changed = sorted(
            n for n in a.namelist() if n in b.namelist() and a.read(n) != b.read(n)
        )
        # the copy's chart + notes must be NEW parts (independence)
        assert any("chart" in n for n in added), added
        assert any("notesSlide" in n for n in added), added
        assert any(n.startswith("ppt/slides/slide") for n in added), added
        # only the wiring parts change; NO existing slide/chart is mutated
        assert changed == [
            "[Content_Types].xml",
            "ppt/_rels/presentation.xml.rels",
            "ppt/presentation.xml",
        ], changed

    def test_duplicated_chart_is_independent(self):
        """Editing the copy's chart does not touch the original's bytes."""
        data = build_three_slide_deck()
        raw = open_raw(data)
        raw.duplicate_slide(1)
        out = raw.to_bytes()
        pkg = zipfile.ZipFile(io.BytesIO(out))
        charts = sorted(n for n in pkg.namelist() if "charts/chart" in n and n.endswith(".xml"))
        assert len(charts) == 2  # original + independent copy
        assert pkg.read(charts[0]) != b"" and pkg.read(charts[1]) != b""

    def test_duplicated_image_is_shared_not_copied(self):
        """The picture on the source slide is REFERENCED by the copy, not
        duplicated (read-only asset)."""
        data = build_deck()  # slide 2 has table + chart + picture
        raw = open_raw(data)
        before = [n for n in zipfile.ZipFile(io.BytesIO(data)).namelist() if "media/image" in n]
        raw.duplicate_slide(1)
        out = raw.to_bytes()
        after = [n for n in zipfile.ZipFile(io.BytesIO(out)).namelist() if "media/image" in n]
        assert len(after) == len(before)  # image shared, not cloned

    def test_duplicate_reopens_in_pptx(self):
        raw = open_raw(build_deck())
        raw.duplicate_slide(0)
        prs = Presentation(io.BytesIO(raw.to_bytes()))
        assert len(prs.slides._sldIdLst) == 3

    def test_out_of_range(self):
        raw = open_raw(build_three_slide_deck())
        with pytest.raises(IndexError):
            raw.duplicate_slide(9)
